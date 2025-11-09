"""
Módulo para cargar y procesar documentos
"""
import os
from typing import Optional
import PyPDF2
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
from io import BytesIO


class DocumentLoader:
    """Carga y extrae texto de diferentes tipos de documentos"""

    # Extensiones de archivos de código soportadas
    CODE_EXTENSIONS = [
        '.py', '.js', '.java', '.cpp', '.c', '.cs', '.php', '.rb', '.go',
        '.rs', '.swift', '.kt', '.ts', '.jsx', '.tsx', '.vue', '.css',
        '.html', '.xml', '.json', '.yaml', '.yml', '.sql', '.sh', '.bat',
        '.r', '.m', '.scala', '.pl', '.lua', '.dart', '.elm'
    ]

    @staticmethod
    def load_pdf(file) -> str:
        """
        Extrae texto de un archivo PDF

        Args:
            file: Archivo PDF cargado

        Returns:
            str: Texto extraído del PDF
        """
        try:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error al leer PDF: {str(e)}")

    @staticmethod
    def load_txt(file) -> str:
        """
        Lee un archivo de texto plano

        Args:
            file: Archivo TXT cargado

        Returns:
            str: Contenido del archivo
        """
        try:
            content = file.read()
            if isinstance(content, bytes):
                # Intentar diferentes encodings
                for encoding in ['utf-8', 'latin-1', 'cp1252']:
                    try:
                        return content.decode(encoding).strip()
                    except UnicodeDecodeError:
                        continue
                # Si ninguno funciona, usar utf-8 con errores ignorados
                return content.decode('utf-8', errors='ignore').strip()
            return content.strip()
        except Exception as e:
            raise Exception(f"Error al leer archivo de texto: {str(e)}")

    @staticmethod
    def load_docx(file) -> str:
        """
        Extrae texto de un archivo DOCX

        Args:
            file: Archivo DOCX cargado

        Returns:
            str: Texto extraído del documento
        """
        try:
            doc = docx.Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error al leer DOCX: {str(e)}")

    @staticmethod
    def load_pptx(file) -> str:
        """
        Extrae texto de un archivo PowerPoint (PPTX)

        Args:
            file: Archivo PPTX cargado

        Returns:
            str: Texto extraído de todas las diapositivas
        """
        try:
            prs = Presentation(file)
            text = ""

            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Diapositiva {slide_num} ---\n"

                # Extraer texto de todas las formas
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text += shape.text + "\n"

                    # Extraer texto de tablas
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text += cell.text + " | "
                            text += "\n"

                    # Extraer notas de las diapositivas
                    if hasattr(slide, "notes_slide") and slide.notes_slide:
                        notes_frame = slide.notes_slide.notes_text_frame
                        if notes_frame and notes_frame.text.strip():
                            text += f"\nNotas: {notes_frame.text}\n"

            return text.strip()
        except Exception as e:
            raise Exception(f"Error al leer PPTX: {str(e)}")

    @staticmethod
    def load_html(file) -> str:
        """
        Extrae texto de un archivo HTML

        Args:
            file: Archivo HTML cargado

        Returns:
            str: Texto extraído del HTML (sin etiquetas)
        """
        try:
            content = file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='ignore')

            # Usar BeautifulSoup para extraer texto
            soup = BeautifulSoup(content, 'lxml')

            # Eliminar scripts y estilos
            for script in soup(["script", "style"]):
                script.decompose()

            # Obtener texto
            text = soup.get_text()

            # Limpiar espacios en blanco excesivos
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)

            return text.strip()
        except Exception as e:
            raise Exception(f"Error al leer HTML: {str(e)}")

    @staticmethod
    def load_code(file, filename: str) -> str:
        """
        Lee un archivo de código fuente

        Args:
            file: Archivo de código cargado
            filename: Nombre del archivo

        Returns:
            str: Contenido del código con metadatos
        """
        try:
            ext = os.path.splitext(filename)[1].lower()
            content = file.read()

            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='ignore')

            # Agregar metadatos sobre el tipo de archivo
            language_map = {
                '.py': 'Python',
                '.js': 'JavaScript',
                '.ts': 'TypeScript',
                '.java': 'Java',
                '.cpp': 'C++',
                '.c': 'C',
                '.cs': 'C#',
                '.php': 'PHP',
                '.rb': 'Ruby',
                '.go': 'Go',
                '.rs': 'Rust',
                '.swift': 'Swift',
                '.kt': 'Kotlin',
                '.r': 'R',
                '.sql': 'SQL',
            }

            language = language_map.get(ext, 'Código')

            formatted_content = f"""
Archivo de código: {filename}
Lenguaje: {language}
Extensión: {ext}

--- Contenido del código ---

{content}
"""
            return formatted_content.strip()
        except Exception as e:
            raise Exception(f"Error al leer archivo de código: {str(e)}")

    @staticmethod
    def load_document(file, filename: str) -> str:
        """
        Detecta el tipo de archivo y extrae el texto

        Args:
            file: Archivo cargado
            filename: Nombre del archivo

        Returns:
            str: Texto extraído
        """
        ext = os.path.splitext(filename)[1].lower()

        if ext == '.pdf':
            return DocumentLoader.load_pdf(file)
        elif ext == '.txt':
            return DocumentLoader.load_txt(file)
        elif ext == '.docx':
            return DocumentLoader.load_docx(file)
        elif ext in ['.pptx', '.ppt']:
            return DocumentLoader.load_pptx(file)
        elif ext in ['.html', '.htm']:
            return DocumentLoader.load_html(file)
        elif ext in DocumentLoader.CODE_EXTENSIONS:
            return DocumentLoader.load_code(file, filename)
        else:
            raise ValueError(f"Tipo de archivo no soportado: {ext}. Formatos soportados: PDF, TXT, DOCX, PPTX, HTML y archivos de código")

    @staticmethod
    def validate_document_length(text: str, min_length: int = 50) -> bool:
        """
        Valida que el documento tenga suficiente contenido

        Args:
            text: Texto del documento
            min_length: Longitud mínima requerida (reducida para código)

        Returns:
            bool: True si el documento es válido
        """
        return len(text.strip()) >= min_length

    @staticmethod
    def get_supported_extensions() -> list:
        """
        Obtiene lista de extensiones soportadas

        Returns:
            list: Lista de extensiones
        """
        return ['pdf', 'txt', 'docx', 'pptx', 'ppt', 'html', 'htm'] + [ext[1:] for ext in DocumentLoader.CODE_EXTENSIONS]
